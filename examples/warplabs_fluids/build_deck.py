"""
Build a 7-slide NVIDIA-style deck for warplabs-fluids WENO5-Z apples-to-apples results.

Run from examples/warplabs_fluids/:
  python build_deck.py
Outputs: warplabs_fluids_deck.pptx
"""

from pathlib import Path
import csv, math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE  = Path(__file__).parent
BENCH = HERE / "benchmarks"
SOD   = BENCH / "sod"
SHO   = BENCH / "shu_osher"
OUT   = HERE / "warplabs_fluids_deck.pptx"

GREEN  = RGBColor(0x76, 0xB9, 0x00)
ORANGE = RGBColor(0xE0, 0x7B, 0x00)
BLACK  = RGBColor(0x00, 0x00, 0x00)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1A, 0x1A, 0x1A)
PANEL  = RGBColor(0x2A, 0x2A, 0x2A)
GRAY   = RGBColor(0xA0, 0xA0, 0xA0)
LGRAY  = RGBColor(0x32, 0x32, 0x32)

SW, SH = Inches(13.333), Inches(7.5)

# ── helpers ───────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def fill_bg(s, color=DARK):
    bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = color

def rect(s, l, t, w, h, bg=None, border=None):
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.line.fill.background()
    if bg:
        sh.fill.solid(); sh.fill.fore_color.rgb = bg
    else:
        sh.fill.background()
    if border:
        sh.line.color.rgb = border; sh.line.width = Pt(1)
    return sh

def txt(s, text, l, t, w, h, size=16, bold=False, italic=False,
        color=WHITE, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb

def ml(s, lines, l, t, w, h, ds=14):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            text, size, color, bold = item, ds, WHITE, False
        else:
            text = item[0]
            size  = item[1] if len(item) > 1 else ds
            color = item[2] if len(item) > 2 else WHITE
            bold  = item[3] if len(item) > 3 else False
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color

def green_bar(s):
    rect(s, 0, 0, 13.333, 0.08, bg=GREEN)

def bottom_bar(s):
    rect(s, 0, 7.3, 13.333, 0.2, bg=PANEL)
    txt(s, "NVIDIA  ·  Warp  ·  warplabs-fluids  |  Phase 1",
        0.2, 7.32, 9, 0.18, size=8, color=GRAY)
    txt(s, "RTX 5000 Ada  ·  Warp 1.12.1  ·  JAX 0.6.2",
        9.8, 7.32, 3.3, 0.18, size=8, color=GRAY, align=PP_ALIGN.RIGHT)

def img(s, path, l, t, w, h=None):
    p = Path(path)
    if not p.exists():
        return
    if h is not None:
        s.shapes.add_picture(str(p), Inches(l), Inches(t),
                             width=Inches(w), height=Inches(h))
    else:
        s.shapes.add_picture(str(p), Inches(l), Inches(t), width=Inches(w))

def chip(s, text, l, t, w=2.8, color=GREEN, tc=BLACK):
    rect(s, l, t, w, 0.28, bg=color)
    txt(s, text, l+0.08, t+0.02, w-0.1, 0.26, size=10, bold=True, color=tc)

def slide_header(s, title, chip_text, chip_w=3.0, chip2=None, chip2_x=None, chip2_w=None):
    green_bar(s)
    txt(s, title, 0.4, 0.12, 12.5, 0.5, size=26, bold=True, color=WHITE)
    chip(s, chip_text, 0.4, 0.72, w=chip_w)
    if chip2 and chip2_x:
        chip(s, chip2, chip2_x, 0.72, w=chip2_w or 2.0, color=ORANGE)
    bottom_bar(s)

# ── CSV readers ───────────────────────────────────────────────────────────────

def read_tp(csv_path):
    d = {}
    with open(csv_path, newline='', encoding='utf-8') as f:
        for row in csv.DictReader(f):
            d.setdefault(row['solver'], {})[int(row['N'])] = float(row['throughput_Mcells'])
    return d

def speedup_at(csv_path, N=4096):
    d = read_tp(csv_path)
    warp = next(v for k, v in d.items() if 'Warp' in k)
    jf   = next(v for k, v in d.items() if 'Jax'  in k)
    return warp[N] / jf[N]

def read_acc(csv_path):
    rows = {}
    with open(csv_path, newline='', encoding='utf-8') as f:
        for row in csv.DictReader(f):
            rows[row['solver']] = row
    return rows

# ── live numbers ──────────────────────────────────────────────────────────────

sod_ratio = speedup_at(SOD / "bench_jaxfluids_throughput.csv")
sho_ratio = speedup_at(SHO / "bench_jaxfluids_throughput.csv")
acc       = read_acc(SOD / "bench_jaxfluids_accuracy.csv")
jf_l1 = float(next(v['L1_rho'] for k, v in acc.items() if 'Jax'  in k))
wp_l1 = float(next(v['L1_rho'] for k, v in acc.items() if 'Warp' in k))
acc_gap = abs(wp_l1 - jf_l1) / jf_l1 * 100

# ── Slide 1 — Title ───────────────────────────────────────────────────────────

def slide_title(prs):
    s = blank_slide(prs); fill_bg(s); green_bar(s)

    rect(s, 0, 0.08, 0.55, 7.22, bg=GREEN)
    txt(s, "WARP", 0.04, 2.9, 0.48, 1.0, size=22, bold=True,
        color=BLACK, align=PP_ALIGN.CENTER)

    txt(s, "NVIDIA Warp  vs  JaxFluids",
        0.75, 1.2, 12.0, 1.0, size=44, bold=True, color=WHITE)
    txt(s, "Apples-to-Apples  ·  WENO5-Z + HLLC + SSP-RK3  ·  1-D Compressible Euler",
        0.75, 2.4, 12.0, 0.5, size=20, color=GREEN)
    txt(s, "Sod Shock Tube  ·  Shu-Osher Density-Wave  ·  Phase 1",
        0.75, 2.98, 12.0, 0.4, size=14, color=GRAY)

    rect(s, 0.75, 3.5, 11.5, 0.03, bg=GRAY)

    stats = [
        (f"{sod_ratio:.0f}×",      "faster than JaxFluids",   "Sod  ·  N=4096 CUDA"),
        (f"{sho_ratio:.0f}×",      "faster than JaxFluids",   "Shu-Osher  ·  N=4096 CUDA"),
        (f"< {math.ceil(acc_gap)}%", "L1(ρ) accuracy gap", "f32 vs f64  ·  N=512"),
        ("f32",                          "Warp precision",          "JaxFluids uses f64"),
    ]
    bw = 2.7; gap = 0.25; x0 = 0.75
    for i, (num, lbl, sub) in enumerate(stats):
        x = x0 + i * (bw + gap)
        rect(s, x, 3.62, bw, 1.5, bg=PANEL)
        rect(s, x, 3.62, bw, 0.05, bg=GREEN)
        txt(s, num,  x+0.14, 3.71, bw-0.28, 0.65, size=34, bold=True, color=GREEN)
        txt(s, lbl,  x+0.14, 4.38, bw-0.28, 0.35, size=12, color=WHITE)
        txt(s, sub,  x+0.14, 4.75, bw-0.28, 0.3,  size=10, color=GRAY)

    txt(s, "RTX 5000 Ada  ·  WSL2 Ubuntu 22.04  ·  Warp 1.12.1  ·  JAX 0.6.2",
        0.75, 6.95, 11.0, 0.3, size=9, color=GRAY)
    bottom_bar(s)


# ── Slide 2 — Algorithm Parity ────────────────────────────────────────────────

def slide_parity(prs):
    s = blank_slide(prs); fill_bg(s)
    slide_header(s, "Algorithm Parity  ·  Closing the Gap",
                 "IDENTICAL ALGORITHM (EXCEPT PRECISION)", chip_w=4.5)

    cols = [
        ("Warp  —  Before",  GRAY,   PANEL),
        ("Warp  —  Now",     GREEN,  RGBColor(0x0F, 0x1E, 0x00)),
        ("JaxFluids",             ORANGE, PANEL),
    ]
    col_x = [0.35, 4.7, 9.05]; col_w = 4.1

    for i, (title, tc, bg) in enumerate(cols):
        rect(s, col_x[i], 1.08, col_w, 0.46, bg=bg, border=tc)
        txt(s, title, col_x[i]+0.15, 1.13, col_w-0.3, 0.38, size=16, bold=True, color=tc)

    rows = [
        ("Reconstruction",   "WENO3  (Jiang-Shu 1996)",  "WENO5-Z  (Borges 2008)",  "WENO5-Z  (Borges 2008)"),
        ("Riemann Solver",   "HLLC",                     "HLLC",                    "HLLC"),
        ("Time Integration", "SSP-RK2  (2 stages)",      "SSP-RK3  (3 stages)",     "SSP-RK3  (3 stages)"),
        ("Ghost cells",      "ng = 2",                   "ng = 3",                  "ng = 3"),
        ("Precision",        "float32",                  "float32",                 "float64"),
        ("Stencil width",    "5 cells",                  "7 cells",                 "7 cells"),
        ("Kernel launches",  "2 / step",                 "3 / step",                "N/A  (XLA JIT)"),
    ]

    for r, (rname, v0, v1, v2) in enumerate(rows):
        y = 1.62 + r * 0.72
        row_bg = PANEL if r % 2 == 0 else LGRAY
        for ci, (cx, val) in enumerate(zip(col_x, [v0, v1, v2])):
            bg = RGBColor(0x0F, 0x1E, 0x00) if ci == 1 else row_bg
            rect(s, cx, y, col_w, 0.68, bg=bg)
            txt(s, rname if ci == 0 else "",
                cx+0.12, y+0.04, col_w-0.24, 0.26, size=11, bold=True, color=GRAY)
            txt(s, val, cx+0.12, y+0.32, col_w-0.24, 0.32,
                size=12, color=(GREEN if ci == 1 else WHITE))

    rect(s, 0.35, 6.78, 12.6, 0.38, bg=GREEN)
    txt(s,
        f"Result:  f32 vs f64 accuracy gap = {acc_gap:.1f}%"
        f"   (Warp {wp_l1:.2e}  vs  JaxFluids {jf_l1:.2e}  at N=512)",
        0.55, 6.82, 12.2, 0.32, size=13, bold=True, color=BLACK)


# ── Slide 3 — Sod Accuracy ────────────────────────────────────────────────────
# Image aspect ratios:
#   jaxfluids_profiles.png  13×4    → 3.25:1  → at w=12.65: h=3.89
#   sod_animation.gif       13×4.2  → 3.10:1  → at w=6.25 : h=2.02
#   sod_convergence.png     15×5    → 3.00:1  → at w=6.10 : h=2.03

def slide_sod(prs):
    s = blank_slide(prs); fill_bg(s)
    slide_header(s, "Sod Shock Tube  ·  Accuracy & Convergence",
                 "WENO5-Z vs EXACT RIEMANN", chip_w=3.2,
                 chip2="ANIMATION", chip2_x=9.9, chip2_w=1.5)

    # profiles — full width, top
    img(s, SOD / "jaxfluids_profiles.png",  0.35, 1.08, 12.65)      # h≈3.89, ends≈4.97

    # animation (left) + convergence (right) — bottom row
    img(s, SOD / "sod_animation.gif",        0.35, 5.05,  6.25)      # h≈2.02, ends≈7.07
    img(s, SOD / "sod_convergence.png",      6.75, 5.05,  6.10)      # h≈2.03, ends≈7.08


# ── Slide 4 — Shu-Osher Accuracy ──────────────────────────────────────────────
# Image aspect ratios:
#   jaxfluids_profiles.png    14×4.5  → 3.11:1  → at w=12.65: h=4.07
#   shu_osher_animation.gif   14×4.5  → 3.11:1  → at w=6.25 : h=2.01
#   shu_osher_convergence.png 14×5    → 2.80:1  → at w=5.70 : h=2.04

def slide_shu_osher(prs):
    s = blank_slide(prs); fill_bg(s)
    slide_header(s, "Shu-Osher Problem  ·  Accuracy & Convergence",
                 "WENO5-Z SELF-CONVERGENCE", chip_w=3.2,
                 chip2="ANIMATION", chip2_x=9.9, chip2_w=1.5)

    img(s, SHO / "jaxfluids_profiles.png",       0.35, 1.08, 12.65)  # h≈4.07, ends≈5.15

    img(s, SHO / "shu_osher_animation.gif",       0.35, 5.25,  6.25) # h≈2.01, ends≈7.26
    img(s, SHO / "shu_osher_convergence.png",     6.75, 5.25,  5.70) # h≈2.04, ends≈7.29


# ── Slide 5 — Throughput (JaxFluids head-to-head) ────────────────────────────
# Image aspect ratios: both 9×6 → 1.5:1  → at w=6.25: h=4.17

def slide_throughput(prs):
    s = blank_slide(prs); fill_bg(s)
    slide_header(s,
        "Performance  ·  Throughput  (Warp f32  vs  JaxFluids f64)",
        "SOD", chip_w=1.1,
        chip2="SHU-OSHER", chip2_x=6.9, chip2_w=1.8)

    img(s, SOD / "jaxfluids_throughput.png",  0.35, 1.08, 6.25)      # h≈4.17, ends≈5.25
    img(s, SHO / "jaxfluids_throughput.png",  6.75, 1.08, 6.25)      # h≈4.17, ends≈5.25

    # callout strip
    rect(s, 0.35, 5.4, 12.6, 1.72, bg=PANEL)
    callouts = [
        (f"{sod_ratio:.0f}×",   "Warp vs JaxFluids\nSod  ·  N=4096"),
        (f"{sho_ratio:.0f}×",   "Warp vs JaxFluids\nShu-Osher  ·  N=4096"),
        ("∼13 Mcell/s",          "Warp CUDA peak\nboth test cases"),
        (f"< {math.ceil(acc_gap)}%", "L1(ρ) accuracy gap\nf32 vs f64  ·  N=512"),
    ]
    cw = 3.05; cx0 = 0.5
    for i, (num, desc) in enumerate(callouts):
        cx = cx0 + i * (cw + 0.13)
        if i > 0:
            rect(s, cx - 0.06, 5.5, 0.01, 1.5, bg=GRAY)
        txt(s, num,  cx, 5.48, cw, 0.6,  size=28, bold=True, color=GREEN)
        txt(s, desc, cx, 6.1,  cw, 0.9,  size=10, color=GRAY)


# ── Slide 6 — GPU Scaling ─────────────────────────────────────────────────────
# Image aspect ratios: both 9×6 → 1.5:1  → at w=6.25: h=4.17

def slide_scaling(prs):
    s = blank_slide(prs); fill_bg(s)
    slide_header(s, "GPU Throughput Scaling  (Warp CPU  ·  Warp CUDA)",
                 "SOD", chip_w=1.1,
                 chip2="SHU-OSHER", chip2_x=6.9, chip2_w=1.8)

    img(s, SOD / "sod_scaling.png",          0.35, 1.08, 6.25)       # h≈4.17, ends≈5.25
    img(s, SHO / "shu_osher_scaling.png",    6.75, 1.08, 6.25)

    ml(s, [
        ("Both cases: Warp CUDA ~767–814 Mcell/s at N=131072", 12, GREEN, True),
        ("Warp CUDA > Warp CPU crossover: N~512  |  Both still bandwidth-scaling at N=131K", 11, GRAY, False),
    ], 0.35, 5.42, 12.6, 0.8, ds=12)


# ── Slide 7 — Memory Footprint ────────────────────────────────────────────────
# Image aspect ratios: both 9×6 → 1.5:1  → at w=6.25: h=4.17

def slide_memory(prs):
    s = blank_slide(prs); fill_bg(s)
    slide_header(s, "GPU Memory Footprint  (Warp  vs  JaxFluids)",
                 "SOD", chip_w=1.1,
                 chip2="SHU-OSHER", chip2_x=6.9, chip2_w=1.8)

    img(s, SOD / "sod_memory.png",           0.35, 1.08, 6.25)
    img(s, SHO / "shu_osher_memory.png",     6.75, 1.08, 6.25)

    ml(s, [
        ("Warp: 32 MiB flat (cudaMallocAsync pool pre-allocated once — O(1) reuse regardless of N)", 12, GREEN, True),
        ("JaxFluids: grows linearly with N  |  Theory: 3 × 3 × (N+6) × 4 B", 11, GRAY, False),
    ], 0.35, 5.42, 12.6, 0.8, ds=12)


# ── assemble ──────────────────────────────────────────────────────────────────

prs = new_prs()
slide_title(prs)
slide_parity(prs)
slide_sod(prs)
slide_shu_osher(prs)
slide_throughput(prs)
slide_scaling(prs)
slide_memory(prs)

prs.save(str(OUT))
print(f"Saved -> {OUT}")
print(f"  Slides: 7")
print(f"  Sod speedup at N=4096:       {sod_ratio:.0f}x")
print(f"  Shu-Osher speedup at N=4096: {sho_ratio:.0f}x")
print(f"  L1(rho) accuracy gap:        {acc_gap:.1f}%")
