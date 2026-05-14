"""
8-slide deck — warplabs-fluids WENO5-Z apples-to-apples results.
White background · Segoe UI · larger type.

python build_deck.py   →   warplabs_fluids_deck.pptx
"""

from pathlib import Path
import csv, math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE  = Path(__file__).parent
BENCH = HERE / "benchmarks"
SOD   = BENCH / "sod"
SHO   = BENCH / "shu_osher"
OUT   = HERE / "warplabs_fluids_deck.pptx"

# ── palette ───────────────────────────────────────────────────────────────────
GREEN  = RGBColor(0x76, 0xB9, 0x00)
GDARK  = RGBColor(0x3E, 0x6A, 0x00)
GLOW   = RGBColor(0xE8, 0xF5, 0xCC)
ORANGE = RGBColor(0xD4, 0x6B, 0x00)
OLIGHT = RGBColor(0xFB, 0xEF, 0xDD)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
INK    = RGBColor(0x12, 0x12, 0x12)
DKGRAY = RGBColor(0x3A, 0x3A, 0x3A)
GRAY   = RGBColor(0x80, 0x80, 0x80)
LGRAY  = RGBColor(0xCC, 0xCC, 0xCC)
PANEL  = RGBColor(0xF4, 0xF5, 0xF6)
PANEL2 = RGBColor(0xE6, 0xE8, 0xEA)

FONT = "Segoe UI"
SW, SH = Inches(13.333), Inches(7.5)

# ── primitives ────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width = SW; prs.slide_height = SH
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def fill_bg(s, color=WHITE):
    bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = color

def rect(s, l, t, w, h, bg=None, border=None, bw=1.0):
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if bg:
        sh.fill.solid(); sh.fill.fore_color.rgb = bg
    else:
        sh.fill.background()
    if border:
        sh.line.color.rgb = border; sh.line.width = Pt(bw)
    else:
        sh.line.fill.background()
    return sh

def txt(s, text, l, t, w, h, sz=18, bold=False, italic=False,
        color=INK, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(sz)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb

def ml(s, lines, l, t, w, h, ds=15, dc=INK):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for item in lines:
        text = item[0] if not isinstance(item, str) else item
        sz   = item[1] if not isinstance(item, str) and len(item) > 1 else ds
        c    = item[2] if not isinstance(item, str) and len(item) > 2 else dc
        bold = item[3] if not isinstance(item, str) and len(item) > 3 else False
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run(); r.text = text
        r.font.name = FONT; r.font.size = Pt(sz)
        r.font.bold = bold; r.font.color.rgb = c

def img(s, path, l, t, w, h=None):
    p = Path(path)
    if not p.exists(): return
    kw = {"width": Inches(w)}
    if h is not None: kw["height"] = Inches(h)
    s.shapes.add_picture(str(p), Inches(l), Inches(t), **kw)

def hline(s, y, l=0.4, w=12.55, color=LGRAY, thickness=0.015):
    rect(s, l, y, w, thickness, bg=color)

def green_bar(s, h=0.12):
    rect(s, 0, 0, 13.333, h, bg=GREEN)

def bottom_bar(s):
    rect(s, 0, 7.28, 13.333, 0.22, bg=PANEL2)
    txt(s, "NVIDIA  ·  Warp  ·  warplabs-fluids  ·  Phase 1",
        0.35, 7.30, 9, 0.20, sz=10, color=GRAY)
    txt(s, "RTX 5000 Ada  ·  Warp 1.12.1  ·  JAX 0.6.2",
        9.5, 7.30, 3.6, 0.20, sz=10, color=GRAY, align=PP_ALIGN.RIGHT)

def slide_header(s, title):
    fill_bg(s); green_bar(s); bottom_bar(s)
    txt(s, title, 0.42, 0.16, 12.5, 0.64, sz=34, bold=True, color=INK)
    hline(s, 0.90)

def tag(s, text, l, t, w, bg=GREEN, tc=WHITE):
    rect(s, l, t, w, 0.32, bg=bg)
    txt(s, text, l+0.12, t+0.05, w-0.24, 0.26, sz=12, bold=True, color=tc)

# ── CSV helpers ───────────────────────────────────────────────────────────────

def read_tp(path):
    d = {}
    with open(path, newline='', encoding='utf-8') as f:
        for row in csv.DictReader(f):
            d.setdefault(row['solver'], {})[int(row['N'])] = float(row['throughput_Mcells'])
    return d

def speedup(path, N=4096):
    d = read_tp(path)
    w = next(v for k, v in d.items() if 'Warp' in k)
    j = next(v for k, v in d.items() if 'Jax'  in k)
    return w[N] / j[N]

def read_acc(path):
    with open(path, newline='', encoding='utf-8') as f:
        return {r['solver']: r for r in csv.DictReader(f)}

# ── live numbers ──────────────────────────────────────────────────────────────

sod_x = speedup(SOD / "bench_jaxfluids_throughput.csv")
sho_x = speedup(SHO / "bench_jaxfluids_throughput.csv")
acc   = read_acc(SOD / "bench_jaxfluids_accuracy.csv")
jf_l1 = float(next(v['L1_rho'] for k, v in acc.items() if 'Jax'  in k))
wp_l1 = float(next(v['L1_rho'] for k, v in acc.items() if 'Warp' in k))
gap   = abs(wp_l1 - jf_l1) / jf_l1 * 100


# ══════════════════════════════════════════════════════════════════════════════
# Slide 1 — Title
# ══════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    s = blank_slide(prs); fill_bg(s)
    green_bar(s, h=0.16)

    # thin left accent
    rect(s, 0, 0.16, 0.10, 7.12, bg=GREEN)

    txt(s, "NVIDIA Warp  vs  JaxFluids",
        0.55, 0.78, 12.5, 1.20, sz=58, bold=True, color=INK)
    txt(s, "Apples-to-Apples  ·  WENO5-Z + HLLC + SSP-RK3",
        0.55, 2.08, 12.5, 0.62, sz=26, color=GREEN)
    txt(s, "1-D Compressible Euler  ·  Sod Shock Tube  ·  Shu-Osher Density-Wave",
        0.55, 2.78, 12.5, 0.44, sz=16, color=GRAY)

    hline(s, 3.32, l=0.55, w=12.2, color=LGRAY)

    # stat cards
    stats = [
        (f"{sod_x:.0f}×",           "faster than JaxFluids",    "Sod  ·  N = 4 096"),
        (f"{sho_x:.0f}×",           "faster than JaxFluids",    "Shu-Osher  ·  N = 4 096"),
        (f"< {math.ceil(gap)}%",    "L1(ρ) accuracy gap",       "Warp f32  vs  JaxFluids f64"),
        ("f32",                      "Warp precision",            "JaxFluids runs f64"),
    ]
    cw = 2.88; gp = 0.22; x0 = 0.55
    for i, (num, lbl, sub) in enumerate(stats):
        x = x0 + i * (cw + gp)
        rect(s, x, 3.46, cw, 1.82, bg=PANEL, border=LGRAY)
        rect(s, x, 3.46, cw, 0.08, bg=GREEN)
        txt(s, num,  x+0.20, 3.58, cw-0.40, 0.78, sz=48, bold=True, color=GREEN)
        txt(s, lbl,  x+0.20, 4.40, cw-0.40, 0.40, sz=14, bold=True, color=INK)
        txt(s, sub,  x+0.20, 4.82, cw-0.40, 0.40, sz=12, color=GRAY)

    txt(s, "RTX 5000 Ada  ·  WSL2 Ubuntu 22.04  ·  Warp 1.12.1  ·  JAX 0.6.2",
        0.55, 6.98, 12.5, 0.28, sz=11, color=LGRAY)
    bottom_bar(s)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 2 — Algorithm Parity (2-column: Warp vs JaxFluids)
# ══════════════════════════════════════════════════════════════════════════════

def slide_parity(prs):
    s = blank_slide(prs)
    slide_header(s, "Algorithm Parity  ·  Warp vs JaxFluids")

    # ── column backgrounds ────────────────────────────────────────────────────
    LEFT_X = 0.42;  COL_W = 5.7
    RIGHT_X = 7.22; COL_W2 = 5.7
    ROWS_Y  = 1.62; ROW_H = 0.72; N_ROWS = 6

    rect(s, LEFT_X,  ROWS_Y - 0.48, COL_W,  N_ROWS*ROW_H + 0.52, bg=GLOW,  border=GREEN, bw=1.5)
    rect(s, RIGHT_X, ROWS_Y - 0.48, COL_W2, N_ROWS*ROW_H + 0.52, bg=OLIGHT, border=ORANGE, bw=1.5)

    # column headers
    rect(s, LEFT_X,  ROWS_Y - 0.48, COL_W,  0.52, bg=GREEN)
    rect(s, RIGHT_X, ROWS_Y - 0.48, COL_W2, 0.52, bg=ORANGE)
    txt(s, "Warp  WENO5-Z  ·  f32",    LEFT_X+0.22,  ROWS_Y-0.44, COL_W-0.44,  0.44,
        sz=20, bold=True, color=WHITE)
    txt(s, "JaxFluids  WENO5-Z  ·  f64", RIGHT_X+0.22, ROWS_Y-0.44, COL_W2-0.44, 0.44,
        sz=20, bold=True, color=WHITE)

    # centre "=" badge
    rect(s, 6.24, ROWS_Y + 0.82, 0.86, 0.86, bg=PANEL2, border=LGRAY)
    txt(s, "=", 6.24, ROWS_Y + 0.82, 0.86, 0.86,
        sz=28, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    rows = [
        ("Reconstruction",    "WENO5-Z  (Borges 2008)",     "WENO5-Z  (Borges 2008)"),
        ("Riemann Solver",    "HLLC",                        "HLLC"),
        ("Time Integration",  "SSP-RK3  ·  3 stages",       "SSP-RK3  ·  3 stages"),
        ("Ghost cells",       "ng = 3  ·  7-cell stencil",   "ng = 3  ·  7-cell stencil"),
        ("Precision",         "float32  ✓",                  "float64"),
        ("Implementation",    "Warp fused GPU kernels",      "JAX  /  XLA JIT"),
    ]
    for i, (label, wval, jval) in enumerate(rows):
        y = ROWS_Y + i * ROW_H
        # row label
        txt(s, label, LEFT_X+0.20, y+0.04, COL_W-0.40, 0.30,
            sz=11, bold=True, color=GDARK)
        txt(s, wval,  LEFT_X+0.20, y+0.34, COL_W-0.40, 0.34, sz=15, color=GDARK)
        txt(s, label, RIGHT_X+0.20, y+0.04, COL_W2-0.40, 0.30,
            sz=11, bold=True, color=ORANGE)
        txt(s, jval,  RIGHT_X+0.20, y+0.34, COL_W2-0.40, 0.34, sz=15, color=DKGRAY)
        # divider
        if i > 0:
            hline(s, y, l=LEFT_X+0.12,  w=COL_W-0.24,  color=GLOW)
            hline(s, y, l=RIGHT_X+0.12, w=COL_W2-0.24, color=OLIGHT)

    # insight bar
    rect(s, 0.42, 6.50, 12.55, 0.62, bg=GREEN)
    txt(s,
        f"Only difference:  f32 vs f64   →   L1(ρ) gap = {gap:.1f}%   "
        f"(Warp {wp_l1:.2e}  vs  JaxFluids {jf_l1:.2e}  at N = 512)",
        0.62, 6.60, 12.15, 0.46, sz=15, bold=True, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# Slides 3 & 4 — Accuracy (Sod / Shu-Osher)
# ══════════════════════════════════════════════════════════════════════════════

def slide_sod(prs):
    s = blank_slide(prs)
    slide_header(s, "Sod Shock Tube  ·  Accuracy & Convergence")
    tag(s, "WENO5-Z  vs  EXACT RIEMANN", 0.42, 0.98, w=3.6)
    img(s, SOD / "jaxfluids_profiles.png", 0.38, 1.40, 12.65)
    img(s, SOD / "sod_animation.gif",       0.38, 5.22,  6.25)
    img(s, SOD / "sod_convergence.png",     6.78, 5.22,  6.10)

def slide_shu_osher(prs):
    s = blank_slide(prs)
    slide_header(s, "Shu-Osher Problem  ·  Accuracy & Convergence")
    tag(s, "WENO5-Z  SELF-CONVERGENCE", 0.42, 0.98, w=3.2)
    img(s, SHO / "jaxfluids_profiles.png",   0.38, 1.40, 12.65)
    img(s, SHO / "shu_osher_animation.gif",   0.38, 5.40,  6.25)
    img(s, SHO / "shu_osher_convergence.png", 6.78, 5.40,  5.70)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 5 — Throughput
# ══════════════════════════════════════════════════════════════════════════════

def slide_throughput(prs):
    s = blank_slide(prs)
    slide_header(s, "Performance  ·  Warp f32  vs  JaxFluids f64")
    tag(s, "SOD",        0.42,  0.98, w=1.1)
    tag(s, "SHU-OSHER",  6.92,  0.98, w=1.9)

    img(s, SOD / "jaxfluids_throughput.png", 0.38, 1.40, 6.25)
    img(s, SHO / "jaxfluids_throughput.png", 6.78, 1.40, 6.25)

    rect(s, 0.38, 5.70, 12.60, 1.42, bg=PANEL, border=LGRAY)
    rect(s, 0.38, 5.70, 12.60, 0.07, bg=GREEN)

    callouts = [
        (f"{sod_x:.0f}×",          "faster than JaxFluids\nSod  ·  N = 4 096"),
        (f"{sho_x:.0f}×",          "faster than JaxFluids\nShu-Osher  ·  N = 4 096"),
        ("∼13 Mcell/s",             "Warp CUDA peak\nboth test cases"),
        (f"< {math.ceil(gap)}%",   "L1(ρ) accuracy gap\nf32 vs f64  ·  N = 512"),
    ]
    cw = 3.0; x0 = 0.55
    for i, (num, desc) in enumerate(callouts):
        cx = x0 + i * (cw + 0.18)
        if i > 0:
            rect(s, cx-0.10, 5.82, 0.012, 1.16, bg=LGRAY)
        txt(s, num,  cx, 5.78, cw, 0.66, sz=42, bold=True, color=GREEN)
        txt(s, desc, cx, 6.46, cw, 0.60, sz=12, color=DKGRAY)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 6 — GPU Scaling
# ══════════════════════════════════════════════════════════════════════════════

def slide_scaling(prs):
    s = blank_slide(prs)
    slide_header(s, "GPU Throughput Scaling")
    tag(s, "SOD",        0.42, 0.98, w=1.1)
    tag(s, "SHU-OSHER",  6.92, 0.98, w=1.9)

    img(s, SOD / "sod_scaling.png",        0.38, 1.40, 6.25)
    img(s, SHO / "shu_osher_scaling.png",  6.78, 1.40, 6.25)

    ml(s, [
        ("Warp CUDA: ~767–814 Mcell/s at N = 131 072  —  still bandwidth-scaling, not yet saturated", 15, GREEN, True),
        ("GPU crossover vs CPU: N ≈ 512   ·   Memory footprint: flat 32 MiB regardless of N  (cudaMallocAsync pool)", 14, DKGRAY, False),
    ], 0.38, 5.66, 12.60, 1.0)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 7 — Memory
# ══════════════════════════════════════════════════════════════════════════════

def slide_memory(prs):
    s = blank_slide(prs)
    slide_header(s, "GPU Memory Footprint")
    tag(s, "SOD",        0.42, 0.98, w=1.1)
    tag(s, "SHU-OSHER",  6.92, 0.98, w=1.9)

    img(s, SOD / "sod_memory.png",        0.38, 1.40, 6.25)
    img(s, SHO / "shu_osher_memory.png",  6.78, 1.40, 6.25)

    ml(s, [
        ("Warp: 32 MiB flat  —  cudaMallocAsync pool pre-allocated once, O(1) reuse regardless of N", 15, GREEN, True),
        ("JaxFluids: grows linearly with N   ·   Theory: 3 × 3 × (N + 6) × 4 B   ·   Crossover: N ≈ 1 000", 14, DKGRAY, False),
    ], 0.38, 5.66, 12.60, 1.0)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 8 — Roadmap  (2 × 4 card grid, 8 phases)
# ══════════════════════════════════════════════════════════════════════════════

PHASES = [
    {
        "n": "1", "done": True,
        "title": "1-D Compressible Euler",
        "sub": "COMPLETE",
        "items": [
            f"WENO5-Z + HLLC + SSP-RK3  ·  float32",
            f"Sod & Shu-Osher V&V",
            f"{sod_x:.0f}× faster than JaxFluids (Sod)",
            f"L1 accuracy gap: {gap:.1f}%  (f32 vs f64)",
        ],
        "bg": GLOW, "border": GREEN, "nbg": GREEN, "ntc": WHITE, "tc": GDARK,
    },
    {
        "n": "2", "done": False,
        "title": "2-D Compressible Euler",
        "sub": "NEXT",
        "items": [
            "Strang dimensional splitting  x/2 → y → x/2",
            "Reuse 1-D fused kernels per axis",
            "Kelvin-Helmholtz instability  V&V",
            "2-D N×N GPU throughput benchmark",
        ],
        "bg": OLIGHT, "border": ORANGE, "nbg": ORANGE, "ntc": WHITE, "tc": DKGRAY,
    },
    {
        "n": "3", "done": False,
        "title": "3-D Compressible Euler",
        "sub": "PLANNED",
        "items": [
            "Full 3-D fused kernel — 3 sweep axes",
            "Rayleigh-Taylor instability  V&V",
            "Shock-vortex interaction benchmark",
            "3-D N³ GPU scaling study",
        ],
        "bg": PANEL, "border": LGRAY, "nbg": PANEL2, "ntc": GRAY, "tc": DKGRAY,
    },
    {
        "n": "4", "done": False,
        "title": "Compressible Navier-Stokes",
        "sub": "PLANNED",
        "items": [
            "Viscous + heat conduction fluxes",
            "Taylor-Green vortex  (Re = 1 600)  V&V",
            "DNS kinetic-energy decay validation",
            "Reynolds-number sweep",
        ],
        "bg": PANEL, "border": LGRAY, "nbg": PANEL2, "ntc": GRAY, "tc": DKGRAY,
    },
    {
        "n": "5", "done": False,
        "title": "Higher-Order Numerics",
        "sub": "PLANNED",
        "items": [
            "WENO7, TENO, MP-WENO schemes",
            "Adaptive stencil reconstruction",
            "Convergence across smooth & shocked regions",
            "Scheme comparison benchmark suite",
        ],
        "bg": PANEL, "border": LGRAY, "nbg": PANEL2, "ntc": GRAY, "tc": DKGRAY,
    },
    {
        "n": "6", "done": False,
        "title": "Two-Phase & Interface Methods",
        "sub": "PLANNED",
        "items": [
            "Diffuse interface  (conservative)",
            "Level-set interface sharpening kernel",
            "Bubble collapse  +  Rayleigh-Plesset V&V",
            "Phase-field GPU benchmark",
        ],
        "bg": PANEL, "border": LGRAY, "nbg": PANEL2, "ntc": GRAY, "tc": DKGRAY,
    },
    {
        "n": "7", "done": False,
        "title": "Multi-GPU & Distributed",
        "sub": "PLANNED",
        "items": [
            "Domain decomposition  (MPI / NCCL)",
            "Warp halo-exchange kernels",
            "Weak & strong scaling  (1 → 64 GPUs)",
            "Target: TB/s collective memory bandwidth",
        ],
        "bg": PANEL, "border": LGRAY, "nbg": PANEL2, "ntc": GRAY, "tc": DKGRAY,
    },
    {
        "n": "8", "done": False,
        "title": "Full JaxFluids Warp Backend",
        "sub": "GOAL",
        "items": [
            "Drop-in Python API — identical config files",
            "All JaxFluids cases run on Warp backend",
            "100× throughput target at production scale",
            "Open-source release alongside JaxFluids",
        ],
        "bg": PANEL, "border": GREEN, "nbg": GDARK, "ntc": WHITE, "tc": DKGRAY,
    },
]

def slide_roadmap(prs):
    s = blank_slide(prs); fill_bg(s); green_bar(s); bottom_bar(s)
    txt(s, "Roadmap  ·  Warp Backend for JaxFluids",
        0.42, 0.16, 12.5, 0.64, sz=34, bold=True, color=INK)
    hline(s, 0.90)

    # 2 rows × 4 columns
    CW = 3.02; CH = 2.84; GAP_X = 0.13; GAP_Y = 0.16
    X0 = (13.333 - 4*CW - 3*GAP_X) / 2   # ≈ 0.455
    Y0 = 1.00

    for idx, ph in enumerate(PHASES):
        row = idx // 4; col = idx % 4
        x = X0 + col * (CW + GAP_X)
        y = Y0 + row * (CH + GAP_Y)

        # card
        rect(s, x, y, CW, CH, bg=ph["bg"], border=ph["border"], bw=1.5)

        # large watermark number (bottom-right, very light)
        wm_color = RGBColor(0xD0, 0xEB, 0xA8) if ph["done"] else PANEL2
        txt(s, ph["n"], x + CW - 0.90, y + CH - 0.88, 0.86, 0.84,
            sz=72, bold=True, color=wm_color, align=PP_ALIGN.RIGHT)

        # phase badge
        rect(s, x+0.16, y+0.16, 0.44, 0.44, bg=ph["nbg"])
        txt(s, ph["n"], x+0.16, y+0.16, 0.44, 0.44,
            sz=17, bold=True, color=ph["ntc"], align=PP_ALIGN.CENTER)

        # sub-label  (COMPLETE / NEXT / PLANNED / GOAL)
        sub_color = GREEN if ph["done"] else (ORANGE if ph["sub"]=="NEXT" else
                    (GDARK  if ph["sub"]=="GOAL"  else GRAY))
        txt(s, ph["sub"], x+0.70, y+0.22, CW-0.84, 0.26,
            sz=10, bold=True, color=sub_color)

        # title
        txt(s, ph["title"], x+0.16, y+0.66, CW-0.32, 0.48,
            sz=13, bold=True, color=ph["tc"])

        hline(s, y+1.18, l=x+0.16, w=CW-0.32,
              color=(GREEN if ph["done"] else (ORANGE if ph["sub"]=="NEXT" else LGRAY)))

        # bullets
        for j, item in enumerate(ph["items"]):
            ty = y + 1.28 + j * 0.36
            rect(s, x+0.20, ty+0.12, 0.06, 0.06,
                 bg=(GREEN if ph["done"] else (ORANGE if ph["sub"]=="NEXT" else LGRAY)))
            txt(s, item, x+0.34, ty, CW-0.50, 0.34,
                sz=10, color=ph["tc"])


# ══════════════════════════════════════════════════════════════════════════════
# Assemble
# ══════════════════════════════════════════════════════════════════════════════

prs = new_prs()
slide_title(prs)
slide_parity(prs)
slide_sod(prs)
slide_shu_osher(prs)
slide_throughput(prs)
slide_scaling(prs)
slide_memory(prs)
slide_roadmap(prs)

prs.save(str(OUT))
print(f"Saved -> {OUT}")
print(f"  8 slides  |  Sod {sod_x:.0f}×  Shu-Osher {sho_x:.0f}×  gap {gap:.1f}%")
